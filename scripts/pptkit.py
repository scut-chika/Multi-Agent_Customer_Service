# -*- coding: utf-8 -*-
"""PPT 生成辅助库：主题样式、文本、形状、表格、箭头等。"""
import math

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# ---------------- 主题色 ----------------
NAVY      = RGBColor(0x0E, 0x2A, 0x47)
NAVY_DEEP = RGBColor(0x08, 0x1C, 0x33)
TEAL      = RGBColor(0x0E, 0xA5, 0xA4)
TEAL_DARK = RGBColor(0x0B, 0x7A, 0x79)
BLUE      = RGBColor(0x3B, 0x82, 0xF6)
LBLUE     = RGBColor(0xDB, 0xEA, 0xFE)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHITE  = RGBColor(0xF4, 0xF7, 0xFA)
INK       = RGBColor(0x1F, 0x29, 0x37)
GRAY      = RGBColor(0x64, 0x74, 0x8B)
LGRAY     = RGBColor(0xE5, 0xE9, 0xEF)
AMBER     = RGBColor(0xE6, 0xA2, 0x3C)
RED       = RGBColor(0xE5, 0x4D, 0x42)
GREEN     = RGBColor(0x2E, 0xA0, 0x43)

FONT = "Microsoft YaHei"
W, H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.55)


def new_prs():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    return prs


def blank_slide(prs, bg=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    rect.fill.solid()
    rect.fill.fore_color.rgb = bg
    rect.line.fill.background()
    rect.shadow.inherit = False
    return slide


def set_run(run, text, size=12, bold=False, color=INK, font=FONT, italic=False):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", font)


def add_text(slide, x, y, w, h, paras, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, wrap=True):
    """paras: list of list-of-runs; run = (text, size, bold, color)"""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    first = True
    for runs in paras:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        for spec in runs:
            text, size, bold, color = spec[0], spec[1], spec[2], spec[3]
            font = spec[4] if len(spec) > 4 else FONT
            italic = spec[5] if len(spec) > 5 else False
            r = p.add_run()
            set_run(r, text, size=size, bold=bold, color=color, font=font, italic=italic)
    return box


def para(runs, align=PP_ALIGN.LEFT, space_after=0, line_spacing=None):
    return {"runs": runs, "align": align, "space_after": space_after, "line_spacing": line_spacing}


def rich_text(slide, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP, wrap=True):
    """paras: list of dict from para(); 支持段间距/行距。"""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    first = True
    for pdef in paras:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = pdef["align"]
        p.space_after = Pt(pdef["space_after"])
        if pdef.get("line_spacing"):
            p.line_spacing = pdef["line_spacing"]
        for spec in pdef["runs"]:
            text, size, bold, color = spec[0], spec[1], spec[2], spec[3]
            font = spec[4] if len(spec) > 4 else FONT
            italic = spec[5] if len(spec) > 5 else False
            r = p.add_run()
            set_run(r, text, size=size, bold=bold, color=color, font=font, italic=italic)
    return box


def rect(slide, x, y, w, h, fill=None, line=None, line_w=1.0,
         shape=MSO_SHAPE.RECTANGLE, radius=None, shadow=False):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = shadow
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp


def fill_shape(sp, fill):
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill


def label(slide, x, y, w, h, text, fill, color=WHITE, size=10, bold=True,
          shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5, line=None, align=PP_ALIGN.CENTER):
    sp = rect(slide, x, y, w, h, fill=fill, line=line, shape=shape, radius=radius)
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = Inches(0.01)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    set_run(r, text, size=size, bold=bold, color=color)
    return sp


def arrow(slide, x1, y1, x2, y2, color=GRAY, weight=1.75, dash=None):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    conn.shadow.inherit = False
    ln = conn.line._get_or_add_ln()
    tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
    ln.append(tail)
    if dash:
        d = ln.makeelement(qn("a:prstDash"), {"val": dash})
        ln.append(d)
    return conn


def header(slide, title, subtitle=None, page=None, total=None):
    rect(slide, MARGIN, Inches(0.5), Inches(0.09), Inches(0.52), fill=TEAL)
    add_text(slide, MARGIN + Inches(0.22), Inches(0.44), Inches(10.5), Inches(0.62),
             [[(title, 24, True, NAVY)]])
    if subtitle:
        add_text(slide, MARGIN + Inches(0.22), Inches(1.02), Inches(11.5), Inches(0.3),
                 [[(subtitle, 11, False, GRAY)]])
    rect(slide, MARGIN, Inches(1.36), W - 2 * MARGIN, Pt(1.2), fill=LGRAY)
    if page:
        total_s = f" / {total}" if total else ""
        add_text(slide, W - Inches(1.1), H - Inches(0.42), Inches(0.7), Inches(0.3),
                 [[(f"{page}{total_s}", 10, False, GRAY)]], align=PP_ALIGN.RIGHT)
    add_text(slide, MARGIN, H - Inches(0.42), Inches(5), Inches(0.3),
             [[("智服方舟 · 多智能体客服自主闭环平台", 9, False, GRAY)]])


def add_table(slide, x, y, w, rows_data, col_widths, row_height=0.32,
              header_fill=NAVY, header_color=WHITE, font_size=9,
              header_size=None, zebra=True, align_map=None, cell_margin=0.05):
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, Inches(row_height * n_rows))
    tbl = tbl_shape.table
    tbl.first_row = False
    tbl.horz_banding = False
    total = sum(col_widths)
    for i, cw in enumerate(col_widths):
        tbl.columns[i].width = Emu(int(w * cw / total))
    for r in range(n_rows):
        tbl.rows[r].height = Inches(row_height)
        for c in range(n_cols):
            cell = tbl.cell(r, c)
            cell.margin_left = cell.margin_right = Inches(cell_margin)
            cell.margin_top = cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if (not zebra or r % 2 == 1) else OFFWHITE
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = (align_map or {}).get(c, PP_ALIGN.LEFT)
            text = str(rows_data[r][c])
            r_ = p.add_run()
            if r == 0:
                set_run(r_, text, size=header_size or (font_size + 0.5), bold=True, color=header_color)
            else:
                set_run(r_, text, size=font_size, bold=False, color=INK)
    return tbl_shape


def section_badge(slide, x, y, w, h, text, fill=NAVY, size=12):
    sp = rect(slide, x, y, w, h, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.35)
    tf = sp.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.06)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    set_run(r, text, size=size, bold=True, color=WHITE)
    return sp


def circle(slide, cx, cy, d, fill=None, line=None, line_w=1.0):
    return rect(slide, cx - d / 2, cy - d / 2, d, d, fill=fill, line=line,
                line_w=line_w, shape=MSO_SHAPE.OVAL)


def text_width_in(text, size_pt, ascii_w=0.55, pad_in=0.45):
    """按字号估算文本渲染宽度（英寸），中文字符按 1.0 字号宽。"""
    ratio = 0.0
    for ch in text:
        ratio += 1.0 if ord(ch) > 0x2E80 else ascii_w
    return ratio * size_pt / 72.0 + pad_in


def hexagon(slide, x, y, w, h, fill, line=None, radius=None):
    return rect(slide, x, y, w, h, fill=fill, line=line, shape=MSO_SHAPE.HEXAGON)


def slide_flow(slide, boxes, arrow_color=GRAY, row_heights=None):
    """boxes: list of dict(x,y,w,h,title,desc,fill,color) —— 供特定页面使用"""
    pass

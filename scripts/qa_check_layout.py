# -*- coding: utf-8 -*-
"""PPT 版式检查：估算文本溢出与越界，输出疑似问题清单。"""
import math
import sys

from pptx import Presentation
from pptx.util import Emu


def cjk_ratio(text):
    """估算字符平均宽度（以字号为单位）：中文≈1.0，ASCII≈0.55。"""
    total = 0.0
    for ch in text:
        if ord(ch) > 0x2E80:
            total += 1.0
        else:
            total += 0.55
    return max(total, 0.1)


def estimate_lines(text, font_pt, width_emu):
    width_pt = width_emu / 12700.0
    per_line = max(int(width_pt / (font_pt * 0.98)), 1)
    lines = 0
    for seg in text.split("\n"):
        n = cjk_ratio(seg)
        lines += max(math.ceil(n / per_line), 1)
    return lines


def main(path):
    prs = Presentation(path)
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    issues = []
    for si, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.left is None or shape.top is None or shape.width is None or shape.height is None:
                continue
            right = shape.left + shape.width
            bottom = shape.top + shape.height
            if right > slide_w + Emu(914400 // 20) or bottom > slide_h + Emu(914400 // 20):
                issues.append((si, shape.name or "?", "越界",
                               f"right={right / 914400:.2f}in bottom={bottom / 914400:.2f}in"))
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            total_h = 0.0
            for p in tf.paragraphs:
                runs = p.runs
                if not runs:
                    continue
                size = max((r.font.size.pt if r.font.size else 12) for r in runs)
                text = "".join(r.text for r in runs)
                n = estimate_lines(text, size, shape.width)
                total_h += n * size * 1.35
            box_h = shape.height / 12700.0
            if total_h > box_h * 1.18 + 4:
                issues.append((si, (shape.name or "?")[:18], "疑似文本溢出",
                               f"估计 {total_h:.0f}pt > 框高 {box_h:.0f}pt"))
        # 表格整体高度估算
        for shape in slide.shapes:
            if shape.has_table:
                tbl = shape.table
                est = 0.0
                for r in range(len(tbl.rows)):
                    row_max = 0.0
                    for c in range(len(tbl.columns)):
                        cell = tbl.cell(r, c)
                        size = 9
                        for p in cell.text_frame.paragraphs:
                            for run in p.runs:
                                if run.font.size:
                                    size = max(size, run.font.size.pt)
                        text = cell.text_frame.text
                        n = estimate_lines(text, size, tbl.columns[c].width)
                        row_max = max(row_max, n * size * 1.35)
                    est += max(row_max + 6, 16)
                total = (shape.top + Emu(int(est * 12700))) / 914400.0
                if total > 7.35:
                    issues.append((si, (shape.name or "?")[:18], "表格疑似超高",
                                   f"底部估计 {total:.2f}in"))
    if issues:
        print(f"共发现 {len(issues)} 个疑似问题：")
        for si, name, kind, detail in issues:
            print(f"  第{si}页 [{name}] {kind}: {detail}")
    else:
        print("未发现明显溢出/越界问题。")
    return 1 if issues else 0


if __name__ == "__main__":
    sys.exit(main(sys.argv[1] if len(sys.argv) > 1 else r"C:\project\ppt\智服方舟-初赛方案.pptx"))
